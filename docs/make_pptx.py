"""Build project_presentation.pptx from scratch with python-pptx.

Run after docs/make_figures.py has populated docs/figures/.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


HERE = Path(__file__).resolve().parent
FIG = HERE / "figures"
OUT = HERE / "project_presentation.pptx"


# --- style helpers ---
TITLE_COLOR = RGBColor(0x1E, 0x3A, 0x5F)   # dark navy
BODY_COLOR = RGBColor(0x1F, 0x2A, 0x3A)
ACCENT_COLOR = RGBColor(0x2B, 0x6C, 0xB0)
LIGHT_GRAY = RGBColor(0x6B, 0x72, 0x80)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_16x9(prs):
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H


def add_blank_slide(prs):
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)


def add_title(slide, text, top=Inches(0.4)):
    box = slide.shapes.add_textbox(Inches(0.6), top, SLIDE_W - Inches(1.2), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    return box


def add_subtitle(slide, text, top):
    box = slide.shapes.add_textbox(Inches(0.6), top, SLIDE_W - Inches(1.2), Inches(0.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(16)
    run.font.color.rgb = LIGHT_GRAY
    run.font.italic = True
    return box


def add_bullets(slide, items, left=Inches(0.6), top=Inches(1.7),
                width=Inches(12.1), height=Inches(5.4), size=18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        run = p.add_run()
        run.text = "•  " + item
        run.font.size = Pt(size)
        run.font.color.rgb = BODY_COLOR
        p.space_after = Pt(10)
    return box


def add_paragraph(slide, text, left=Inches(0.6), top=Inches(1.7),
                  width=Inches(12.1), height=Inches(2.0), size=18, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.italic = italic
    run.font.color.rgb = BODY_COLOR
    return box


def add_image(slide, path, left, top, width=None, height=None):
    if not Path(path).exists():
        print(f"  !! missing figure: {path}")
        return None
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def add_footer(slide, text="NDWI Water-Body Segmentation · GNR Semester 8"):
    box = slide.shapes.add_textbox(Inches(0.6), SLIDE_H - Inches(0.5),
                                   SLIDE_W - Inches(1.2), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.color.rgb = LIGHT_GRAY
    run.font.italic = True


# --- slides ---

def slide_title(prs):
    s = add_blank_slide(prs)
    # big title
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.3), SLIDE_W - Inches(1.6), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "NDWI Water-Body Segmentation"
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = TITLE_COLOR

    # subtitle
    box2 = s.shapes.add_textbox(Inches(0.8), Inches(3.6), SLIDE_W - Inches(1.6), Inches(1.0))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "Extracting water bodies from satellite imagery using NDWI + Mean-Shift segmentation"
    r2.font.size = Pt(20)
    r2.font.color.rgb = LIGHT_GRAY

    # team
    box3 = s.shapes.add_textbox(Inches(0.8), Inches(5.3), SLIDE_W - Inches(1.6), Inches(1.5))
    tf3 = box3.text_frame
    for i, line in enumerate([
        "GNR Semester 8 — Problem Statement #6",
        "Tanmay Mandaliya · Rohan · Pravesh Khaparde",
    ]):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(16)
        r.font.color.rgb = BODY_COLOR if i == 1 else ACCENT_COLOR


def slide_problem(prs):
    s = add_blank_slide(prs)
    add_title(s, "The problem")
    add_subtitle(s, "What the assignment asks for", Inches(1.25))
    add_paragraph(s, (
        '"Generate an NDWI image and make a 3-band image of Green, NIR, and '
        'NDWI components. Extract all water bodies from this 3-band image '
        'using the Mean-Shift segmentation algorithm."'
    ), top=Inches(1.9), italic=True, size=18)
    add_bullets(s, [
        "Input: a Green-band and a NIR-band satellite GeoTIFF of the same place.",
        "Output: a binary water mask (white = water, black = everything else) plus area in km² for each water body.",
        "No training data. No labels. Unsupervised clustering only.",
        "Must be defended as our own implementation — no sklearn / cv2 Mean-Shift.",
    ], top=Inches(3.3), size=18)
    add_footer(s)


def slide_ndwi(prs):
    s = add_blank_slide(prs)
    add_title(s, "NDWI — why water pops out")

    # formula box
    form_box = s.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.2))
    tf = form_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "NDWI  =  (Green − NIR) / (Green + NIR)"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = ACCENT_COLOR

    add_bullets(s, [
        "Water reflects Green light but absorbs NIR → NDWI is large and positive.",
        "Vegetation does the opposite (absorbs Green, reflects NIR) → NDWI is negative.",
        "Built-up and bare soil → NDWI is near zero.",
        "Pure arithmetic, applied pixel-by-pixel. No training, no model.",
    ], top=Inches(3.2), size=18)

    # range strip (visualised as text)
    strip = s.shapes.add_textbox(Inches(0.6), Inches(5.7), Inches(12.1), Inches(0.7))
    tf = strip.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "NDWI range:   −1   ⇐   land / vegetation          built-up ≈ 0          water  ⇒   +1"
    r.font.size = Pt(14)
    r.font.color.rgb = LIGHT_GRAY
    r.font.italic = True
    add_footer(s)


def slide_feature_space(prs):
    s = add_blank_slide(prs)
    add_title(s, "Stacking into a 3-D feature space")
    add_subtitle(s, "Every pixel becomes one point in a 3-D feature space", Inches(1.25))
    add_bullets(s, [
        "Step: stack Green, NIR, NDWI into a single array of shape (H, W, 3).",
        "Every pixel is now a 3-D point with coordinates (Green, NIR, NDWI).",
        "Water pixels cluster together in one region of this 3-D space.",
        "Vegetation pixels cluster in a different region; built-up in yet another.",
        "Clustering this point cloud is the job of Mean-Shift.",
    ], top=Inches(2.0), size=18)
    add_footer(s)


def slide_meanshift_intuition(prs):
    s = add_blank_slide(prs)
    add_title(s, "Mean-Shift — the core idea")
    add_image(s, FIG / "meanshift_concept.png",
              left=Inches(0.6), top=Inches(1.4), width=Inches(12.1))
    add_paragraph(s, (
        "Each point shifts toward the weighted mean of its neighbours within "
        "a radius called the bandwidth. Repeating this shrinks the cloud into "
        "a handful of density peaks. Points that converge to the same peak "
        "are one cluster."
    ), top=Inches(6.3), height=Inches(0.9), size=14, italic=True)
    add_footer(s)


def slide_pseudocode(prs):
    s = add_blank_slide(prs)
    add_title(s, "Our segment() algorithm — written from scratch")
    code = (
        "def segment(img3, bandwidth, max_iter=100, eps=1e-3):\n"
        "    points = img3.reshape(H*W, 3)\n"
        "    # normalise so bandwidth is isotropic\n"
        "    pts = (points - points.mean(0)) / (points.std(0) + 1e-9)\n"
        "    tree = cKDTree(pts)         # fast neighbour lookup\n"
        "\n"
        "    for i in range(N):          # shift every pixel-point\n"
        "        x = pts[i]\n"
        "        while True:\n"
        "            nbrs = pts[ tree.query_ball_point(x, bandwidth) ]\n"
        "            w    = exp(-||nbrs - x||² / (2·h²))    # Gaussian weights\n"
        "            x_new = (w · nbrs).sum() / w.sum()      # weighted mean\n"
        "            if |x_new − x| < eps: break\n"
        "            x = x_new\n"
        "        shifted[i] = x\n"
        "\n"
        "    # points landing within bandwidth/2 → same cluster\n"
        "    return merge_modes(shifted).reshape(H, W)"
    )
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = code
    r.font.name = "Consolas"
    r.font.size = Pt(15)
    r.font.color.rgb = BODY_COLOR
    add_footer(s)


def slide_pipeline(prs):
    s = add_blank_slide(prs)
    add_title(s, "The full pipeline, in one picture")
    add_image(s, FIG / "pipeline_tahoe.png",
              left=Inches(0.3), top=Inches(1.35), width=Inches(12.7))
    add_paragraph(s, (
        "Lake Tahoe sample: (1) raw Green band  →  (2) NDWI coloured blue→red  "
        "→  (3) Mean-Shift clusters (random colours)  →  (4) final water mask "
        "overlaid on Green."
    ), top=Inches(6.4), height=Inches(0.8), size=13, italic=True)
    add_footer(s)


def slide_results(prs):
    s = add_blank_slide(prs)
    add_title(s, "Results on other scenes")
    add_image(s, FIG / "result_mumbai.png",
              left=Inches(0.4), top=Inches(1.3), width=Inches(6.2))
    add_image(s, FIG / "result_khadakwasla.png",
              left=Inches(6.8), top=Inches(1.3), width=Inches(6.2))
    add_paragraph(s, "Left: Mumbai coast — captures sea and creek. Right: Khadakwasla reservoir.",
                  top=Inches(5.7), height=Inches(0.5), size=14, italic=True)
    add_bullets(s, [
        "10 Sentinel-2 scenes bundled: Tahoe, Mumbai, Dal Lake, Chilika, Amazon, Sundarbans, Great Salt Lake, Khadakwasla, Ganges, Iceland.",
        "All pulled from AWS public Sentinel-2 COGs — no authentication needed.",
    ], top=Inches(6.3), size=14)
    add_footer(s)


def slide_ui(prs):
    s = add_blank_slide(prs)
    add_title(s, "The Streamlit UI")
    add_bullets(s, [
        "Sidebar: pick a scene from a dropdown, tweak four parameters (downsample k, bandwidth, min-NDWI, min-blob).",
        "Five tabs: Input · NDWI · Segmentation · Water mask · Stats.",
        "Click Run — Mean-Shift takes ~30-60 s with a progress bar.",
        "Stats tab reports: how many water bodies, total area in km², per-blob breakdown.",
        "No file upload, no download buttons — kept simple so the algorithm is the focus.",
    ], top=Inches(1.8), size=18)
    add_footer(s)


def slide_summary(prs):
    s = add_blank_slide(prs)
    add_title(s, "Summary — what is ours vs what is a library")

    # Table-like text
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12.1), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True

    rows = [
        ("Ours (we wrote the code)",
         "NDWI formula · band stacking · Mean-Shift shift loop · Gaussian kernel · convergence check · mode merging · water-cluster selection · blob cleanup"),
        ("Plumbing (library call, labelled as such)",
         "rasterio — GeoTIFF I/O · scipy.cKDTree — fast neighbour search · scipy.ndimage.label — connected components · numpy — array math · streamlit/matplotlib — UI"),
        ("Deliberately NOT used",
         "sklearn.cluster.MeanShift · cv2.pyrMeanShiftFiltering — would defeat the project"),
    ]
    for i, (h, body) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r1 = p.add_run()
        r1.text = h + "\n"
        r1.font.size = Pt(16)
        r1.font.bold = True
        r1.font.color.rgb = ACCENT_COLOR
        r2 = p.add_run()
        r2.text = body
        r2.font.size = Pt(14)
        r2.font.color.rgb = BODY_COLOR
        p.space_after = Pt(14)

    # closing line
    thanks = s.shapes.add_textbox(Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.6))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Thank you — questions welcome."
    r.font.size = Pt(18)
    r.font.italic = True
    r.font.color.rgb = TITLE_COLOR
    add_footer(s)


def main():
    prs = Presentation()
    set_16x9(prs)

    slide_title(prs)
    slide_problem(prs)
    slide_ndwi(prs)
    slide_feature_space(prs)
    slide_meanshift_intuition(prs)
    slide_pseudocode(prs)
    slide_pipeline(prs)
    slide_results(prs)
    slide_ui(prs)
    slide_summary(prs)

    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
